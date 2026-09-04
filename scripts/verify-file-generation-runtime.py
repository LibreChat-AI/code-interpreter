#!/usr/bin/env python3
"""Verify native file generation inside the sandbox runtime."""

from __future__ import annotations

import csv
import json
import math
import platform
import struct
import subprocess
import tempfile
import wave
import xml.etree.ElementTree as ET
import zipfile
from importlib import metadata
from packaging.version import Version
from pathlib import Path

from docx import Document
from openpyxl import Workbook, load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from reportlab.pdfgen import canvas
import tableau_parser
import tableauserverclient


TABLEAU_HYPER_SUPPORTED = platform.machine().lower() in {"amd64", "x86_64"}
if TABLEAU_HYPER_SUPPORTED:
    from tableauhyperapi import (
        Connection,
        CreateMode,
        HyperProcess,
        Inserter,
        SqlType,
        TableDefinition,
        TableName,
        Telemetry,
    )


REQUIRED_DISTRIBUTIONS = (
    "openpyxl",
    "xlsxwriter",
    "pandas",
    "numpy",
    "python-docx",
    "docxtpl",
    "python-pptx",
    "pillow",
    "pypdf",
    "pypdf2",
    "pdfplumber",
    "reportlab",
    "tableauserverclient",
    "tableau-parser",
    "weasyprint",
    "msgpack",
    "setuptools",
)

TABLEAU_DISTRIBUTIONS = {
    "tableauserverclient": "0.41",
    "tableau-parser": "0.1.0",
}
if TABLEAU_HYPER_SUPPORTED:
    TABLEAU_DISTRIBUTIONS["tableauhyperapi"] = "0.0.26359"
    REQUIRED_DISTRIBUTIONS += ("tableauhyperapi",)


def run(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(args, check=True, capture_output=True, text=True)


def ffprobe(path: Path) -> dict:
    result = run(
        "ffprobe",
        "-v",
        "error",
        "-show_streams",
        "-show_format",
        "-of",
        "json",
        str(path),
    )
    return json.loads(result.stdout)


def require(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def require_ooxml(path: Path, member: str) -> None:
    require(zipfile.is_zipfile(path), f"{path.name} is not a ZIP-based OOXML file")
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
    require("[Content_Types].xml" in names, f"{path.name} has no OOXML content types")
    require(member in names, f"{path.name} is missing {member}")


def require_pdf(path: Path) -> None:
    require(path.stat().st_size > 100, f"{path.name} is empty")
    require(path.read_bytes()[:4] == b"%PDF", f"{path.name} has no PDF signature")
    require(len(PdfReader(path).pages) >= 1, f"{path.name} cannot be parsed")


def require_media(
    path: Path,
    *,
    format_name: str,
    codecs: dict[str, str],
    minimum_duration: float = 0.8,
) -> None:
    require(path.stat().st_size > 100, f"{path.name} is empty")
    info = ffprobe(path)
    actual_format = str(info.get("format", {}).get("format_name", ""))
    require(format_name in actual_format.split(","), f"{path.name} container is {actual_format}")

    duration = float(info.get("format", {}).get("duration") or 0)
    require(duration >= minimum_duration, f"{path.name} duration is {duration}")

    streams = {
        str(stream.get("codec_type")): str(stream.get("codec_name"))
        for stream in info.get("streams", [])
    }
    for stream_type, codec in codecs.items():
        require(
            streams.get(stream_type) == codec,
            f"{path.name} {stream_type} codec is {streams.get(stream_type)}",
        )


def command_version(*args: str) -> str:
    result = run(*args)
    output = result.stdout.strip() or result.stderr.strip()
    return output.splitlines()[0]


def runtime_versions() -> dict[str, object]:
    packages = {name: metadata.version(name) for name in REQUIRED_DISTRIBUTIONS}
    for name, expected_version in TABLEAU_DISTRIBUTIONS.items():
        require(
            packages[name] == expected_version,
            f"{name} {packages[name]} does not match required {expected_version}",
        )
    require(
        Version(packages["weasyprint"]) >= Version("68"),
        f"WeasyPrint {packages['weasyprint']} is below the patched runtime baseline",
    )
    require(
        Version(packages["msgpack"]) >= Version("1.2.1"),
        f"msgpack {packages['msgpack']} is below the patched runtime baseline",
    )
    require(
        Version(packages["setuptools"]) >= Version("78.1.1"),
        f"setuptools {packages['setuptools']} is below the patched runtime baseline",
    )
    return {
        "python": platform.python_version(),
        "libreoffice": command_version("libreoffice", "--version"),
        "ffmpeg": command_version("ffmpeg", "-version"),
        "unzip": command_version("unzip", "-v"),
        "weasyprint": command_version("weasyprint", "--version"),
        "packages": packages,
    }


def package_bundle_checksum() -> str:
    checksum_path = Path("/pkgs/.bundle.sha256")
    require(checksum_path.is_file(), "/pkgs/.bundle.sha256 is missing")
    checksum = checksum_path.read_text(encoding="ascii").strip()
    require(len(checksum) == 64, "package bundle checksum has an invalid length")
    require(
        all(character in "0123456789abcdef" for character in checksum),
        "package bundle checksum is invalid",
    )
    return checksum


def main() -> None:
    results: dict[str, str] = {}
    with tempfile.TemporaryDirectory(prefix="wowloop-file-runtime-") as tmp:
        root = Path(tmp)

        twb = root / "workbook.twb"
        twb.write_text(
            '<?xml version="1.0" encoding="utf-8"?><workbook version="2024.1"/>',
            encoding="utf-8",
        )
        twbx = root / "workbook.twbx"
        with zipfile.ZipFile(twbx, "w") as archive:
            archive.write(twb, arcname=twb.name)
        extracted = root / "twbx"
        extracted.mkdir()
        run("unzip", "-q", str(twbx), "-d", str(extracted))
        require(
            extracted.joinpath(twb.name).read_text(encoding="utf-8")
            == twb.read_text(encoding="utf-8"),
            "twbx extraction failed",
        )
        require(tableau_parser.__name__ == "tableau_parser", "tableau-parser import failed")
        require(
            tableauserverclient.__name__ == "tableauserverclient",
            "tableauserverclient import failed",
        )
        results["twbx"] = "pass"

        if TABLEAU_HYPER_SUPPORTED:
            hyper_path = root / "extract.hyper"
            table = TableDefinition(
                TableName("Extract", "Extract"),
                [TableDefinition.Column("name", SqlType.text())],
            )
            with HyperProcess(Telemetry.DO_NOT_SEND_USAGE_DATA_TO_TABLEAU) as hyper:
                with Connection(
                    hyper.endpoint,
                    hyper_path,
                    CreateMode.CREATE_AND_REPLACE,
                ) as connection:
                    connection.catalog.create_schema("Extract")
                    connection.catalog.create_table(table)
                    with Inserter(connection, table) as inserter:
                        inserter.add_row(["WowLoop"])
                        inserter.execute()
                    rows = connection.execute_list_query(
                        'SELECT "name" FROM "Extract"."Extract"'
                    )
            require(rows == [["WowLoop"]], "Hyper extract round trip failed")
            results["hyper"] = "pass"
        else:
            results["hyper"] = f"unsupported_on_{platform.machine().lower()}"

        xlsx = root / "report.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Launch Summary"
        sheet.append(["Metric", "Value"])
        sheet.append(["Revenue", 1250])
        sheet["B3"] = "=B2*1.1"
        workbook.save(xlsx)
        reopened = load_workbook(xlsx, data_only=False)
        require(reopened["Launch Summary"]["B3"].value == "=B2*1.1", "xlsx formula missing")
        require_ooxml(xlsx, "xl/workbook.xml")
        results["xlsx"] = "pass"

        fake_xlsx = root / "renamed-text.xlsx"
        fake_xlsx.write_text("not a workbook", encoding="utf-8")
        try:
            load_workbook(fake_xlsx)
        except zipfile.BadZipFile:
            pass
        else:
            raise AssertionError("renamed text passed XLSX validation")
        results["renamed_text_rejected"] = "pass"

        csv_path = root / "report.csv"
        with csv_path.open("w", encoding="utf-8", newline="") as handle:
            csv.writer(handle).writerows([["name", "value"], ["conversion", "12.5%"]])
        with csv_path.open(encoding="utf-8", newline="") as handle:
            require(list(csv.reader(handle))[1] == ["conversion", "12.5%"], "csv parse failed")
        results["csv"] = "pass"

        docx = root / "report.docx"
        document = Document()
        document.add_heading("WowLoop Report", 0)
        document.add_paragraph("Validated structured document output.")
        document.save(docx)
        require(Document(docx).paragraphs[1].text.startswith("Validated"), "docx reopen failed")
        require_ooxml(docx, "word/document.xml")
        results["docx"] = "pass"

        pptx = root / "report.pptx"
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = "WowLoop"
        slide.placeholders[1].text = "Validated presentation output"
        deck.save(pptx)
        require(Presentation(pptx).slides[0].shapes.title.text == "WowLoop", "pptx reopen failed")
        require_ooxml(pptx, "ppt/presentation.xml")
        results["pptx"] = "pass"

        pdf = root / "report.pdf"
        pdf_canvas = canvas.Canvas(str(pdf))
        pdf_canvas.drawString(72, 760, "WowLoop validated PDF output")
        pdf_canvas.save()
        require_pdf(pdf)
        results["pdf"] = "pass"

        for image_format, suffix in (("PNG", "png"), ("JPEG", "jpg"), ("WEBP", "webp")):
            image_path = root / f"report.{suffix}"
            Image.new("RGB", (160, 90), (28, 99, 220)).save(image_path, format=image_format)
            with Image.open(image_path) as image:
                image.verify()
            with Image.open(image_path) as image:
                require(image.format == image_format, f"{image_format} format mismatch")
                require(image.size == (160, 90), f"{image_format} size mismatch")
            results[suffix] = "pass"

        svg = root / "report.svg"
        svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="160" height="90">'
            '<rect width="160" height="90" fill="#1c63dc"/></svg>',
            encoding="utf-8",
        )
        require(ET.parse(svg).getroot().tag.endswith("svg"), "svg parse failed")
        results["svg"] = "pass"

        wav = root / "report.wav"
        sample_rate = 16000
        with wave.open(str(wav), "wb") as audio:
            audio.setnchannels(1)
            audio.setsampwidth(2)
            audio.setframerate(sample_rate)
            frames = bytearray()
            for index in range(sample_rate):
                sample = int(8000 * math.sin(2 * math.pi * 440 * index / sample_rate))
                frames.extend(struct.pack("<h", sample))
            audio.writeframes(frames)
        with wave.open(str(wav), "rb") as audio:
            require(audio.getnframes() == sample_rate, "wav reopen failed")
        results["wav"] = "pass"

        mp3 = root / "report.mp3"
        run("ffmpeg", "-y", "-i", str(wav), "-codec:a", "libmp3lame", str(mp3))
        require_media(mp3, format_name="mp3", codecs={"audio": "mp3"})
        results["mp3"] = "pass"

        mp4 = root / "report.mp4"
        run(
            "ffmpeg",
            "-y",
            "-f",
            "lavfi",
            "-i",
            "color=c=0x1c63dc:s=320x180:d=1",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=440:duration=1",
            "-shortest",
            "-c:v",
            "libx264",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            str(mp4),
        )
        require_media(mp4, format_name="mp4", codecs={"video": "h264", "audio": "aac"})
        results["mp4"] = "pass"

        office_source = root / "office-source.docx"
        Document(docx).save(office_source)
        run(
            "libreoffice",
            f"-env:UserInstallation={root.joinpath('libreoffice-profile').as_uri()}",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(root),
            str(office_source),
        )
        converted_pdf = root / "office-source.pdf"
        require_pdf(converted_pdf)
        results["libreoffice_conversion"] = "pass"

        html = root / "report.html"
        html.write_text("<h1>WowLoop</h1><p>Validated HTML to PDF conversion.</p>", encoding="utf-8")
        html_pdf = root / "html-report.pdf"
        run("weasyprint", str(html), str(html_pdf))
        require_pdf(html_pdf)
        results["html_to_pdf"] = "pass"

    print(
        json.dumps(
            {
                "status": "pass",
                "formats": results,
                "package_bundle_checksum": package_bundle_checksum(),
                "versions": runtime_versions(),
            },
            sort_keys=True,
        )
    )


if __name__ == "__main__":
    main()
